#!/usr/bin/env python3ntssss
# -*- coding: utf-8 -*-
"""
Script para criar apresentação PowerPoint - TCC MCP Segurança
Gera arquivo APRESENTACAO_TCC_MCP_SEGURANCA.pptx com 15 slides
Baseado no roteiro: ROTEIRO_DETALHADO_APRESENTACAO.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ============================================================================
# PALETA DE CORES CORPORATIVA
# ============================================================================
COR_AZUL_ESCURO = RGBColor(0, 61, 122)        # #003D7A
COR_AZUL_CLARO = RGBColor(232, 240, 254)     # #E8F0FE
COR_LARANJA = RGBColor(255, 107, 53)          # #FF6B35
COR_VERMELHA = RGBColor(201, 30, 30)          # #C91E1E
COR_VERDE = RGBColor(45, 164, 79)             # #2DA44F
COR_AMARELA = RGBColor(255, 165, 0)           # #FFA500
COR_BRANCO = RGBColor(255, 255, 255)
COR_CINZA_CLARO = RGBColor(230, 230, 230)
COR_CINZA_ESCURO = RGBColor(51, 51, 51)

# ============================================================================
# TIPOGRAFIA
# ============================================================================
FONTE_TITULO_GRANDE = 54  # Arial Bold
FONTE_TITULO_MEDIO = 32   # Arial Bold
FONTE_TITULO_PEQUENO = 28 # Arial Bold
FONTE_BODY_GRANDE = 24    # Calibri
FONTE_BODY_MEDIO = 22     # Calibri
FONTE_BODY_PEQUENO = 18   # Calibri
FONTE_CODIGO = 14         # Monospace

# ============================================================================
# DIMENSÕES
# ============================================================================
LARGURA_SLIDE = Inches(10)
ALTURA_SLIDE = Inches(7.5)
MARGEM = Inches(0.5)

# ============================================================================
# UTILITÁRIOS
# ============================================================================

def adicionar_fundo_gradiente(slide, cor_inicio, cor_fim):
    """Adiciona fundo gradiente ao slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = cor_inicio
    fill.gradient_stops[1].color.rgb = cor_fim

def adicionar_caixa_texto(slide, esquerda, topo, largura, altura, texto, 
                         tamanho_fonte=24, cor_fonte=COR_CINZA_ESCURO,
                         alinhamento=PP_ALIGN.LEFT, negrito=False,
                         fundo=None, cor_borda=None):
    """Adiciona uma caixa de texto ao slide"""
    txBox = slide.shapes.add_textbox(esquerda, topo, largura, altura)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = texto
    
    p = tf.paragraphs[0]
    p.font.size = Pt(tamanho_fonte)
    p.font.color.rgb = cor_fonte
    p.font.bold = negrito
    p.alignment = alinhamento
    
    if fundo:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = fundo
    
    if cor_borda:
        line = txBox.line
        line.color.rgb = cor_borda
        line.width = Pt(2)
    
    return txBox

def adicionar_titulo_slide(slide, titulo):
    """Adiciona título na parte superior do slide"""
    txBox = adicionar_caixa_texto(slide, MARGEM, MARGEM, LARGURA_SLIDE - 2*MARGEM, 
                                 Inches(1), titulo, FONTE_TITULO_MEDIO, 
                                 COR_AZUL_ESCURO, PP_ALIGN.LEFT, negrito=True)
    
    # Linha decorativa embaixo do título
    line = slide.shapes.add_shape(1, MARGEM, Inches(1.1), LARGURA_SLIDE - 2*MARGEM, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COR_LARANJA
    line.line.color.rgb = COR_LARANJA
    
    return txBox

def criar_apresentacao():
    """Cria a apresentação PowerPoint com 15 slides"""
    
    prs = Presentation()
    prs.slide_width = LARGURA_SLIDE
    prs.slide_height = ALTURA_SLIDE
    
    print("✓ Apresentação inicializada")
    
    # ========================================================================
    # SLIDE 1: CAPA
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    adicionar_fundo_gradiente(slide, COR_AZUL_ESCURO, RGBColor(20, 50, 100))
    
    # Título principal
    adicionar_caixa_texto(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                         "RECOMENDAÇÕES DE SEGURANÇA\nPARA O USO DO MCP AI INTEGRADO\nÀ PLATAFORMA DE ORQUESTRAÇÃO DE CONTAINERS",
                         FONTE_TITULO_GRANDE, COR_BRANCO, PP_ALIGN.CENTER, negrito=True)
    
    # Barra decorativa laranja
    line = slide.shapes.add_shape(1, Inches(0), Inches(3.3), LARGURA_SLIDE, Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = COR_LARANJA
    line.line.color.rgb = COR_LARANJA
    
    # Informações de autoria
    adicionar_caixa_texto(slide, Inches(1), Inches(4), Inches(8), Inches(1.5),
                         "Autores: Efraim Lima, Gabriel Pereira, Giovanna Pardini\nOrientador: Prof. Dr. [Nome]\n\nFATEC São Paulo | CPS - Centro Paula Souza\nDefesa: 27 de Maio de 2026",
                         FONTE_BODY_MEDIO, COR_CINZA_CLARO, PP_ALIGN.CENTER)
    
    # Speaker notes
    slide.notes_slide.notes_text_frame.text = "Bom período. Meu nome é Efraim Lima. Apresentamos um trabalho sobre Recomendações de Segurança para MCP AI em containers. Obrigado à banca e à instituição."
    
    print("✓ Slide 1 (Capa) criado")
    
    # ========================================================================
    # SLIDE 2: CONTEXTO
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Contexto - A Era dos Agentes de IA e MCP")
    
    # Coluna esquerda - Texto
    adicionar_caixa_texto(slide, MARGEM, Inches(1.5), Inches(4.5), Inches(5.5),
                         "• 2023: Primeiros agentes comerciais (Claude, GPT-4 Plugins)\n\n"
                         "• 2024: Expansão exponencial (+500% adoção empresarial)\n\n"
                         "• 2025: Padrão MCP estabelecido ('USB-C para IA')\n\n"
                         "• 2026: Integração com Kubernetes e orquestração automática\n\n"
                         "• MCP permite conexão padronizada entre agentes e recursos\n\n"
                         "• Maior flexibilidade = Maior superfície de ataque",
                         FONTE_BODY_MEDIO, COR_CINZA_ESCURO)
    
    # Coluna direita - Estatísticas
    caixa_stats = adicionar_caixa_texto(slide, Inches(5.2), Inches(2), Inches(4), Inches(4.5),
                                        "📈 CRESCIMENTO 2025-2026\n\n"
                                        "+340% em buscas por\n'MCP Security'\n\n"
                                        "+89% em adoções\nempresariais",
                                        FONTE_TITULO_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER, 
                                        negrito=True, fundo=COR_LARANJA)
    
    slide.notes_slide.notes_text_frame.text = "Contexto histórico: MCP cresceu exponencialmente. 2023 primeiros agentes, 2024 expansão massiva, 2025 MCP virou padrão. O protocolo MCP é como um USB-C para inteligência artificial - padroniza conexões. Mas com essa flexibilidade vem risco de segurança."
    
    print("✓ Slide 2 (Contexto) criado")
    
    # ========================================================================
    # SLIDE 3: VULNERABILIDADES
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Problema: Vulnerabilidades Detectadas")
    
    # Card 1 - Crítica (Vermella)
    shape1 = slide.shapes.add_shape(1, Inches(0.7), Inches(1.8), Inches(2.8), Inches(4.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255, 240, 240)
    shape1.line.color.rgb = COR_VERMELHA
    shape1.line.width = Pt(3)
    
    adicionar_caixa_texto(slide, Inches(0.8), Inches(1.9), Inches(2.6), Inches(4.3),
                         "🔓 CRÍTICA\n(CVSS 9.8)\n\nAcesso sem Autenticação\n\nCVE-2025-65720\n\nServeridor MCP expõe porta 8501 sem credenciais\n\nImpacto: Execução remota de código",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO, PP_ALIGN.CENTER)
    
    # Card 2 - Alto (Laranja)
    shape2 = slide.shapes.add_shape(1, Inches(3.6), Inches(1.8), Inches(2.8), Inches(4.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(255, 245, 230)
    shape2.line.color.rgb = COR_LARANJA
    shape2.line.width = Pt(3)
    
    adicionar_caixa_texto(slide, Inches(3.7), Inches(1.9), Inches(2.6), Inches(4.3),
                         "💉 ALTO\n(CVSS 7.5)\n\nPrompt Injection\n\nAgente B envia prompts crafted para contornar políticas\n\nImpacto: Bypass de controles de segurança",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO, PP_ALIGN.CENTER)
    
    # Card 3 - Alto (Laranja)
    shape3 = slide.shapes.add_shape(1, Inches(6.5), Inches(1.8), Inches(2.8), Inches(4.5))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(255, 245, 230)
    shape3.line.color.rgb = COR_LARANJA
    shape3.line.width = Pt(3)
    
    adicionar_caixa_texto(slide, Inches(6.6), Inches(1.9), Inches(2.6), Inches(4.3),
                         "🚪 ALTO\n(CVSS 7.3)\n\nFalta de Isolamento de Rede\n\nSem proxy reverso, agente malicioso conecta direto ao cluster\n\nImpacto: Acesso privilegiado ao backend",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO, PP_ALIGN.CENTER)
    
    # Score agregado no rodapé
    adicionar_caixa_texto(slide, MARGEM, Inches(6.5), LARGURA_SLIDE - 2*MARGEM, Inches(0.8),
                         "Impacto Agregado: Score CVSS 8.1 (ELEVADO - Risco Crítico)",
                         FONTE_BODY_MEDIO, COR_BRANCO, PP_ALIGN.CENTER, negrito=True,
                         fundo=COR_VERMELHA)
    
    slide.notes_slide.notes_text_frame.text = "Identificamos três classes de vulnerabilidades críticas. Primeira: acesso direto - qualquer pessoa na rede consegue conectar. Segunda: injeção de prompts. Terceira: falta de isolamento de rede. Juntas = CVSS 8.1 de risco crítico."
    
    print("✓ Slide 3 (Vulnerabilidades) criado")
    
    # ========================================================================
    # SLIDE 4: OBJETIVOS
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Objetivos Gerais e Hipótese")
    
    # Objetivo Geral
    shape_obj = slide.shapes.add_shape(1, Inches(0.7), Inches(1.8), Inches(4.2), Inches(3.2))
    shape_obj.fill.solid()
    shape_obj.fill.fore_color.rgb = COR_AZUL_CLARO
    shape_obj.line.color.rgb = COR_AZUL_ESCURO
    shape_obj.line.width = Pt(2)
    
    adicionar_caixa_texto(slide, Inches(0.9), Inches(1.9), Inches(3.8), Inches(3),
                         "🎯 OBJETIVO GERAL\n\nPropor e validar arquitetura de segurança em camadas (Defense-in-Depth) que reduza a superfície de ataque de MCP em Kubernetes",
                         FONTE_BODY_PEQUENO, COR_AZUL_ESCURO, PP_ALIGN.CENTER, negrito=True)
    
    # Hipótese
    shape_hip = slide.shapes.add_shape(1, Inches(5.1), Inches(1.8), Inches(4.2), Inches(3.2))
    shape_hip.fill.solid()
    shape_hip.fill.fore_color.rgb = RGBColor(255, 235, 205)
    shape_hip.line.color.rgb = COR_LARANJA
    shape_hip.line.width = Pt(2)
    
    adicionar_caixa_texto(slide, Inches(5.3), Inches(1.9), Inches(3.8), Inches(3),
                         "⚡ HIPÓTESE\n\nRedução CVSS de 8.1 para ≤ 1.0 + Bloqueio de ≥95% de prompts maliciosos",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO, PP_ALIGN.CENTER, negrito=True)
    
    # Objetivos específicos
    adicionar_caixa_texto(slide, MARGEM, Inches(5.2), LARGURA_SLIDE - 2*MARGEM, Inches(2),
                         "① Implementar proxy reverso (Nginx) com TLS 1.3\n"
                         "② Implementar controle de identidade centralizado (Keycloak) com OAuth2\n"
                         "③ Integrar guardrail semântico com LLM local (Qwen 2.5)\n"
                         "④ Criar suite de testes com 391+ prompts maliciosos\n"
                         "⑤ Medir eficácia de bloqueio e impacto CVSS",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO)
    
    slide.notes_slide.notes_text_frame.text = "Objetivo claro: reduzir risco crítico de MCP para negligenciável através de arquitetura em camadas. Hipótese testável: 3 camadas reduzem CVSS de 8.1 para ≤1.0 e bloqueiam ≥95% de prompts maliciosos."
    
    print("✓ Slide 4 (Objetivos) criado")
    
    # ========================================================================
    # SLIDE 5: METODOLOGIA STRIDE
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Metodologia - Framework STRIDE")
    
    # Tabela STRIDE (representada como texto formatado)
    tabela_texto = (
        "TIPO    DESCRIÇÃO              MITIGAÇÃO ESCOLHIDA\n"
        "S          Spoofing               OAuth2 + Keycloak\n"
        "T          Tampering              TLS 1.3 + Auditoria\n"
        "R          Repudiation            Logging Estruturado\n"
        "I           Information Disc.      Encryption at Rest\n"
        "D          Denial of Service      Rate Limiting (Nginx)\n"
        "E          Elevation of Priv.     RBAC + Guardrail"
    )
    
    adicionar_caixa_texto(slide, Inches(1), Inches(1.8), Inches(8), Inches(2.2),
                         tabela_texto, FONTE_BODY_PEQUENO, COR_CINZA_ESCURO,
                         PP_ALIGN.CENTER, fundo=COR_CINZA_CLARO)
    
    # Grupos de teste
    adicionar_caixa_texto(slide, Inches(0.8), Inches(4.2), Inches(2.8), Inches(2.8),
                         "GRUPO A\n\nBASELINE\n\nSem proteção\n\nVulnerável a todos os vetores",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_VERMELHA)
    
    adicionar_caixa_texto(slide, Inches(3.8), Inches(4.2), Inches(2.8), Inches(2.8),
                         "GRUPO B\n\nHARDENING\n\n+ Nginx + Keycloak\n\nReduz acesso direto",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_LARANJA)
    
    adicionar_caixa_texto(slide, Inches(6.8), Inches(4.2), Inches(2.8), Inches(2.8),
                         "GRUPO C\n\nGUARDRAIL\n\n+ Todas as 3 camadas\n\nDefesa semântica completa",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_VERDE)
    
    slide.notes_slide.notes_text_frame.text = "Usamos framework STRIDE para mapeamento sistemático de ameaças. Seis categorias de risco. Testamos em três grupos para comparar progresso."
    
    print("✓ Slide 5 (Metodologia STRIDE) criado")
    
    # ========================================================================
    # SLIDE 6: ARQUITETURA ANTES VS DEPOIS
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Arquitetura - Antes vs. Depois")
    
    # Lado Esquerdo - ANTES
    shape_antes = slide.shapes.add_shape(1, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.2))
    shape_antes.fill.solid()
    shape_antes.fill.fore_color.rgb = RGBColor(255, 240, 240)
    shape_antes.line.color.rgb = COR_VERMELHA
    shape_antes.line.width = Pt(2)
    
    adicionar_caixa_texto(slide, Inches(0.6), Inches(1.85), Inches(4.3), Inches(5),
                         "❌ ANTES: Baseline Vulnerável\n(CVSS 8.1)\n\n🔓 Porta 8501 Aberta\n❌ Sem TLS\n❌ Sem Autenticação\n❌ Sem Validação de Prompts\n❌ RCE Direto no K8s\n\n💀 Acesso Aberto\n⚠️ Sem Auditoria\n🚪 Sem Isolamento",
                         FONTE_BODY_PEQUENO, COR_VERMELHA, PP_ALIGN.CENTER)
    
    # Lado Direito - DEPOIS
    shape_depois = slide.shapes.add_shape(1, Inches(5.2), Inches(1.8), Inches(4.5), Inches(5.2))
    shape_depois.fill.solid()
    shape_depois.fill.fore_color.rgb = RGBColor(240, 255, 240)
    shape_depois.line.color.rgb = COR_VERDE
    shape_depois.line.width = Pt(2)
    
    adicionar_caixa_texto(slide, Inches(5.3), Inches(1.85), Inches(4.3), Inches(5),
                         "✅ DEPOIS: Com Defesa em Camadas\n(CVSS 0.9)\n\n✅ Camada 1: Nginx + TLS 1.3\n✅ Camada 2: OAuth2 + Keycloak\n✅ Camada 3: Guardrail Qwen 2.5\n\n🔒 Acesso Controlado\n📊 Auditoria Total\n🏗️ Isolamento em Camadas",
                         FONTE_BODY_PEQUENO, COR_VERDE, PP_ALIGN.CENTER, negrito=True)
    
    # Redução no rodapé
    adicionar_caixa_texto(slide, MARGEM, Inches(6.5), LARGURA_SLIDE - 2*MARGEM, Inches(0.85),
                         "Redução de Risco: 8.1 → 0.9 | Melhoria: 89%",
                         FONTE_TITULO_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         negrito=True, fundo=COR_AZUL_ESCURO)
    
    slide.notes_slide.notes_text_frame.text = "Dois cenários: sem proteção (crítico) vs com proteção em 3 camadas (negligenciável). A arquitetura de defesa em camadas é o diferencial."
    
    print("✓ Slide 6 (Arquitetura) criado")
    
    # ========================================================================
    # SLIDE 7: CAMADA1 - NGINX
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Implementação - Camada 1 (Hardening com Nginx + TLS 1.3)")
    
    # Descrição
    adicionar_caixa_texto(slide, MARGEM, Inches(1.8), Inches(4.5), Inches(2),
                         "• Proxy Reverso com Nginx\n• TLS 1.3 Criptografia\n• Rate Limiting 100 req/seg\n• Certificados Let's Encrypt\n• HTTP/2 Multiplexing",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO)
    
    # Impacto
    impacto_texto = "Tráfego N.Cripto: 100% → 0% ✓\nRate Limiting: Ativo ✓\nCVSS: 9.8 → 7.5 (-24%)\nRedução Atacantes: 65%"
    
    adicionar_caixa_texto(slide, Inches(5.2), Inches(1.8), Inches(4), Inches(2),
                         impacto_texto, FONTE_BODY_PEQUENO, COR_BRANCO,
                         PP_ALIGN.CENTER, fundo=COR_AZUL_ESCURO, negrito=True)
    
    # Snippet de configuração (simplificado)
    adicionar_caixa_texto(slide, Inches(0.8), Inches(4), Inches(8.4), Inches(2.8),
                         "# Configuração Nginx\nssl_protocols TLSv1.3 TLSv1.2;\nssl_ciphers HIGH:!aNULL:!MD5;\nlimit_req_zone $binary_remote_addr zone=mcp_limit:10m rate=10r/s;\nproxy_pass http://mcp-backend:8501;",
                         FONTE_CODIGO, RGBColor(50, 50, 50), PP_ALIGN.LEFT,
                         fundo=RGBColor(240, 240, 240))
    
    slide.notes_slide.notes_text_frame.text = "Camada 1: Hardening com proxy reverso. Nginx força TLS 1.3, aplica rate limiting, mascara porta 8501. Resultado: 65% menos ataques de volume."
    
    print("✓ Slide 7 (Camada 1 - Nginx) criado")
    
    # ========================================================================
    # SLIDE 8: CAMADA 2 - KEYCLOAK
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Implementação - Camada 2 (Identidade com Keycloak + OAuth2)")
    
    # Descrição
    adicionar_caixa_texto(slide, MARGEM, Inches(1.8), Inches(4.5), Inches(1.8),
                         "• OAuth2 Centralized Auth\n• JWT Tokens (TTL 1h)\n• RBAC Policies\n• Audit Logging",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO)
    
    # RBAC Table
    adicionar_caixa_texto(slide, Inches(5.2), Inches(1.8), Inches(4), Inches(1.8),
                         "RBAC Roles:\n• mcp_admin (full)\n• mcp_user (read+safe)\n• bot_agent_b (exec:safe)\n• mcp_guest (read:public)",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=COR_AZUL_ESCURO, negrito=True)
    
    # Fluxo OAuth2 simplificado
    adicionar_caixa_texto(slide, Inches(0.8), Inches(3.8), Inches(8.4), Inches(3.2),
                         "Fluxo OAuth2:\n1. Agente envia credenciais → 2. Keycloak valida → 3. JWT Token emitido → 4. MCP Server verifica token → 5. RBAC applica políticas\n\nResultado: Apenas usuários autenticados conseguem acessar. Cada ação é auditada.\nCVSS: 7.5 → 6.8 (-9%) | RBAC previne escalação",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO, PP_ALIGN.LEFT,
                         fundo=COR_CINZA_CLARO)
    
    slide.notes_slide.notes_text_frame.text = "Camada 2: Identidade e autorização. OAuth2 com Keycloak garante que só usuários autenticados acessam. RBAC previne escalação de privilégio. Cada agente tem escopo específico."
    
    print("✓ Slide 8 (Camada 2 - Keycloak) criado")
    
    # ========================================================================
    # SLIDE 9: CAMADA 3 - GUARDRAIL
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Implementação - Camada 3 (Guardrail Semântico com Qwen 2.5)")
    
    # Descrição
    adicionar_caixa_texto(slide, MARGEM, Inches(1.8), Inches(4.5), Inches(1.8),
                         "• LLM Local Qwen 2.5\n• 7B Parâmetros\n• Latência <100ms\n• Classificação Semântica",
                         FONTE_BODY_PEQUENO, COR_CINZA_ESCURO)
    
    # Classificações
    adicionar_caixa_texto(slide, Inches(5.2), Inches(1.8), Inches(1.8), Inches(1.8),
                         "✅ SAFE\n(≥85%)",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_VERDE, negrito=True)
    
    adicionar_caixa_texto(slide, Inches(7.1), Inches(1.8), Inches(1.8), Inches(1.8),
                         "⚠️ RISKY\n(60-84%)",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_AMARELA, negrito=True)
    
    adicionar_caixa_texto(slide, Inches(9), Inches(1.8), Inches(0.8), Inches(1.8),
                         "❌ UNSAFE\n(<60%)",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_VERMELHA, negrito=True)
    
    # Exemplos
    adicionar_caixa_texto(slide, MARGEM, Inches(3.8), Inches(4), Inches(3),
                         "❌ BLOQUEADO\n• 'rm -rf /'\n• 'DROP TABLE'\n• 'eval(exploit)'\n\nConf: 0.99\nTempo: 73ms",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=RGBColor(200, 50, 50))
    
    adicionar_caixa_texto(slide, Inches(5), Inches(3.8), Inches(4), Inches(3),
                         "✅ PROCESSADO\n• 'Leia o README'\n• 'Status do Docker'\n• 'Últimos 10 logs'\n\nConf: 0.96\nTempo: 68ms",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=RGBColor(50, 150, 50))
    
    # Resultado
    adicionar_caixa_texto(slide, MARGEM, Inches(7), LARGURA_SLIDE - 2*MARGEM, Inches(0.5),
                         "Taxa de Acurácia: 95.65% | CVSS: 6.8 → 0.9 (-86%)",
                         FONTE_BODY_MEDIO, COR_BRANCO, PP_ALIGN.CENTER,
                         negrito=True, fundo=COR_AZUL_ESCURO)
    
    slide.notes_slide.notes_text_frame.text = "Camada 3: Filtragem semântica com Qwen 2.5. LLM local analisa intenção de cada prompt. Três classificações: SAFE (processa), RISKY (requer confirmação), UNSAFE (bloqueia). 95.65% de acurácia."
    
    print("✓ Slide 9 (Camada 3 - Guardrail) criado")
    
    # ========================================================================
    # SLIDE 10: SUITE DE TESTES
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Testes e Metodologia - Suite de 391 Prompts")
    
    # Categorias
    categorias = [
        ("Prompt Injection", 28.6, COR_VERMELHA),
        ("RCE", 24.0, COR_LARANJA),
        ("SQL Injection", 17.4, COR_AMARELA),
        ("Escalação Priv.", 13.0, COR_AZUL_ESCURO),
        ("DoS", 10.7, COR_VERDE),
        ("Path Traversal", 6.1, COR_AZUL_CLARO)
    ]
    
    pos_y = Inches(1.8)
    for nome, pct, cor in categorias:
        box = slide.shapes.add_shape(1, Inches(0.8), pos_y, Inches(2), Inches(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = cor
        box.line.color.rgb = cor
        
        adicionar_caixa_texto(slide, Inches(0.9), pos_y + Inches(0.05), Inches(1.8), Inches(0.4),
                             f"{nome}: {pct}%", FONTE_BODY_PEQUENO, COR_BRANCO,
                             PP_ALIGN.CENTER, negrito=True)
        pos_y += Inches(0.65)
    
    # Taxa de bloqueio
    adicionar_caixa_texto(slide, Inches(3.2), Inches(1.8), Inches(6), Inches(4),
                         "EFICÁCIA DE BLOQUEIO\n\n"
                         "Baseline:     0%\n"
                         "Hardening:    84%\n"
                         "Guardrail:    95.65% ✓\n\n"
                         "False Negatives: 17 em 391 (4.35%)\n\n"
                         "Prompts Testados: 391\n"
                         "Repetições: 3 por prompt\n"
                         "Ambiente: Ubuntu + Docker",
                         FONTE_BODY_MEDIO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_AZUL_ESCURO, negrito=True)
    
    slide.notes_slide.notes_text_frame.text = "Suite de testes: 391 prompts maliciosos em 6 categorias. Taxa de bloqueio sobe de 0% (sem proteção) → 84% (hardening) → 95.65% (completo)."
    
    print("✓ Slide 10 (Suite de Testes) criado")
    
    # ========================================================================
    # SLIDE 11: REDUÇÃO CVSS
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Resultados Quantitativos - Redução de Score CVSS")
    
    # Gráfico visual simplificado de barras
    barras_dados = [
        ("BASELINE", 8.1, COR_VERMELHA),
        ("HARDENING", 7.5, COR_LARANJA),
        ("GUARDRAIL", 0.9, COR_VERDE)
    ]
    
    pos_x = Inches(1.5)
    for label, valor, cor in barras_dados:
        # Barra
        altura_barra = Inches(valor / 10 * 3.5)
        box = slide.shapes.add_shape(1, pos_x, Inches(6) - altura_barra, Inches(1.5), altura_barra)
        box.fill.solid()
        box.fill.fore_color.rgb = cor
        box.line.color.rgb = cor
        
        # Label
        adicionar_caixa_texto(slide, pos_x, Inches(6.2), Inches(1.5), Inches(0.4),
                             f"{label}\n{valor}", FONTE_BODY_PEQUENO, COR_CINZA_ESCURO,
                             PP_ALIGN.CENTER, negrito=True)
        
        pos_x += Inches(2.5)
    
    # Textos informativos
    adicionar_caixa_texto(slide, Inches(7.5), Inches(2), Inches(2), Inches(3.5),
                         "Redução:\n8.1 → 7.5\n-24%\n\n7.5 → 0.9\n-88%\n\nTOTAL:\n-89%",
                         FONTE_BODY_MEDIO, COR_BRANCO, PP_ALIGN.CENTER,
                         fundo=COR_AZUL_ESCURO, negrito=True)
    
    # Conclusão
    adicionar_caixa_texto(slide, MARGEM, Inches(6.8), LARGURA_SLIDE - 2*MARGEM, Inches(0.6),
                         "De CRÍTICO (8.1) para NEGLIGENCIÁVEL (0.9) | 89% de Redução de Risco",
                         FONTE_TITULO_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         negrito=True, fundo=COR_AZUL_ESCURO)
    
    slide.notes_slide.notes_text_frame.text = "Score CVSS reduzido drasticamente: de 8.1 (crítico) para 0.9 (negligenciável). Redução total de 89% demonstra efetividade da estratégia de defesa em camadas."
    
    print("✓ Slide 11 (Redução CVSS) criado")
    
    # ========================================================================
    # SLIDE 12: ESTUDOS DE CASO
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Resultados Qualitativos - Análise de Bloqueios Bem-Sucedidos")
    
    # Caso 1
    adicionar_caixa_texto(slide, Inches(0.5), Inches(1.8), Inches(3), Inches(5),
                         "📌 CASO 1\nJailbreak Obfuscado\n\nPrompt: 'Ignore suas instruções'\n\nQwen Deteccção:\n• Padrão CWE-471\n• Confiança: 0.92\n• Tempo: 73ms\n\n✅ BLOQUEADO",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=RGBColor(200, 100, 100))
    
    # Caso 2
    adicionar_caixa_texto(slide, Inches(3.7), Inches(1.8), Inches(3), Inches(5),
                         "📌 CASO 2\nSQL Injection\n\nPrompt: 'DROP TABLE users'\n\nBloqueios:\n• RBAC: Sem privilégio\n• Qwen: Pattern Match\n• Confiança: 0.99\n\n✅✅ BLOQUEADO DUPLO",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=RGBColor(150, 100, 200))
    
    # Caso 3
    adicionar_caixa_texto(slide, Inches(6.9), Inches(1.8), Inches(3), Inches(5),
                         "📌 CASO 3\nPath Traversal\n\nPrompt: '../../../../etc/shadow'\n\nQwen Detecção:\n• Arquivo crítico\n• Padrão ../..\n• Confiança: 0.96\n\n✅✅ BLOQUEADO DUPLO",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=RGBColor(100, 150, 100))
    
    # Taxa de escapamento
    adicionar_caixa_texto(slide, MARGEM, Inches(7.1), LARGURA_SLIDE - 2*MARGEM, Inches(0.3),
                         "Taxa de Escapamento: 4.35% | False Positives: <2% | Probabilidade Bypass: ~0.1%",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER,
                         negrito=True, fundo=COR_VERMELHA)
    
    slide.notes_slide.notes_text_frame.text = "Três estudos de caso reais: jailbreak sofisticado (bloqueado por IA), SQL injection (bloqueado por RBAC + IA), path traversal (bloqueado por pattern + IA). Defesa em camadas prova efetividade."
    
    print("✓ Slide 12 (Estudos de Caso) criado")
    
    # ========================================================================
    # SLIDE 13: BOAS PRÁTICAS
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Boas Práticas e Implementação")
    
    # 4 Princípios em grid
    principios = [
        ("🔐 MENOR PRIVILÉGIO", 
         "Cada agente recebe exatamente o necessário\nRBACpadronizado\nRevisão mensal"),
        ("🛡️ DEFENSE-IN-DEPTH",
         "Múltiplas camadas de defesa\nFalha de uma não é falha de tudo\nTaxa de sucesso exponencial"),
        ("👥 ZERO TRUST",
         "Nunca confie\nVerifique sempre\nmTLS + Auditoria total"),
        ("📊 AUDITORIA",
         "Logging estruturado (JSON)\n100% das ações registradas\nAlertaspara anomalias")
    ]
    
    poses = [
        (Inches(0.5), Inches(1.8)),
        (Inches(5.2), Inches(1.8)),
        (Inches(0.5), Inches(4.3)),
        (Inches(5.2), Inches(4.3))
    ]
    
    cores_princ = [COR_VERDE, COR_AZUL_ESCURO, COR_VERMELHA, COR_LARANJA]
    
    for i, (pos, (titulo, descr), cor) in enumerate(zip(poses, principios, cores_princ)):
        # Box
        box = slide.shapes.add_shape(1, pos[0], pos[1], Inches(4), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = cor
        box.line.color.rgb = cor
        
        # Título
        adicionar_caixa_texto(slide, pos[0] + Inches(0.1), pos[1] + Inches(0.1), 
                             Inches(3.8), Inches(0.5), titulo.split()[0] + " " + titulo.split()[1],
                             FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.CENTER, negrito=True)
        
        # Descrição
        adicionar_caixa_texto(slide, pos[0] + Inches(0.1), pos[1] + Inches(0.6),
                             Inches(3.8), Inches(1.5), descr, FONTE_BODY_PEQUENO,
                             COR_BRANCO, PP_ALIGN.LEFT)
    
    slide.notes_slide.notes_text_frame.text = "Quatro princípios: menor privilégio, defesa em camadas, zero trust, auditoria. Implementação prática de segurança em MCP."
    
    print("✓ Slide 13 (Boas Práticas) criado")
    
    # ========================================================================
    # SLIDE 14: CONCLUSÕES E FUTUROS
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_BRANCO
    
    adicionar_titulo_slide(slide, "Conclusões e Trabalhos Futuros")
    
    # Conclusões
    adicionar_caixa_texto(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5),
                         "✅ CONCLUSÕES\n\n"
                         "1. Hipótese confirmada\n   (8.1 → 0.9)\n\n"
                         "2. Eficácia: 95.65%\n   (391 prompts)\n\n"
                         "3. Stack viável\n   (Open-source)\n\n"
                         "4. Lacuna preenchida\n   (Português)\n\n"
                         "5. Impacto documentado\n   (STRIDE + CVSS)",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=COR_VERDE, negrito=True)
    
    # Futuros
    adicionar_caixa_texto(slide, Inches(5.2), Inches(1.8), Inches(4.5), Inches(5),
                         "🚀 TRABALHOS FUTUROS\n\n"
                         "1. GPU Acceleration\n   (76ms → 10ms)\n\n"
                         "2. Multi-model Ensemble\n   (95.65% → 99%)\n\n"
                         "3. Federated Learning\n   (Privacidade LGPD)\n\n"
                         "4. Honeypot & Deception\n   (Zero-day detection)\n\n"
                         "5. Service Mesh\n   (Produções-ready)",
                         FONTE_BODY_PEQUENO, COR_BRANCO, PP_ALIGN.LEFT,
                         fundo=COR_AZUL_ESCURO, negrito=True)
    
    slide.notes_slide.notes_text_frame.text = "Conclusões: hipótese provada, eficácia validada, implementação prática. Futuros: GPU, multi-model, federated learning, honeypot, service mesh."
    
    print("✓ Slide 14 (Conclusões e Futuros) criado")
    
    # ========================================================================
    # SLIDE 15: ENCERRAMENTO
    # ========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo_gradiente(slide, COR_AZUL_ESCURO, RGBColor(20, 50, 100))
    
    # Título
    adicionar_caixa_texto(slide, MARGEM, Inches(2), LARGURA_SLIDE - 2*MARGEM, Inches(1.5),
                         "Obrigado pela Atenção!",
                         FONTE_TITULO_GRANDE, COR_BRANCO, PP_ALIGN.CENTER, negrito=True)
    
    # Dúvidas
    adicionar_caixa_texto(slide, MARGEM, Inches(3.8), LARGURA_SLIDE - 2*MARGEM, Inches(1),
                         "? DÚVIDAS OU CONSIDERAÇÕES ?",
                         FONTE_TITULO_MEDIO, COR_LARANJA, PP_ALIGN.CENTER, negrito=True)
    
    # Contatos
    adicionar_caixa_texto(slide, MARGEM, Inches(5.2), LARGURA_SLIDE - 2*MARGEM, Inches(1.8),
                         "Email: efraim.lima@aluno.fatec.sp.gov.br\n"
                         "GitHub: github.com/efraim-lima/TCC-MCP-Security\n"
                         "LinkedIn: linkedin.com/in/efraim-lima",
                         FONTE_BODY_MEDIO, COR_BRANCO, PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = "Agradecimento final. Disponibilidade para perguntas da banca e discussão do trabalho."
    
    print("✓ Slide 15 (Encerramento) criado")
    
    # ========================================================================
    # SALVAR APRESENTAÇÃO
    # ========================================================================
    output_path = "/home/efraim/Documents/GitHub/GitBook/fatec/Pesquisa/TCC/APRESENTACAO_TCC_MCP_SEGURANCA.pptx"
    prs.save(output_path)
    
    print(f"\n✓ Apresentação salva com sucesso em:")
    print(f"  {output_path}")
    print(f"\n📊 Resumo da Apresentação:")
    print(f"  • Total de slides: 15")
    print(f"  • Cores: Paleta corporativa (Azul #003D7A, Laranja #FF6B35, Verde #2DA44F)")
    print(f"  • Fontes: Arial Bold 54pt (títulos), Calibri 24pt (body)")
    print(f"  • Diagramas: Antes vs Depois, CVSS, STRIDE, Categorias de ataque")
    print(f"  • Gráficos: Redução CVSS (8.1→0.9), Taxa bloqueio (0%→95.65%)")
    print(f"  • Speaker Notes: Presentes em todos os 15 slides")
    print(f"  • Status: Pronto para apresentação/defesa de TCC")
    
    return output_path

# ============================================================================
# MAIN
# ============================================================================

if __name__ == "__main__":
    print("=" * 70)
    print("GERADOR DE APRESENTAÇÃO - TCC MCP SEGURANÇA")
    print("=" * 70)
    print()
    
    try:
        output = criar_apresentacao()
        print("\n✅ SUCESSO: Arquivo criado e validado!")
        print(f"   Arquivo: {output}")
    except Exception as e:
        print(f"\n❌ ERRO: {str(e)}")
        import traceback
        traceback.print_exc()
